"""
Generate PPT for JPX Stock Prediction Project Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
TITLE_COLOR = RGBColor(0, 51, 102)  # Dark blue
ACCENT_COLOR = RGBColor(0, 102, 204)  # Blue


def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(12.333), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)

    return slide


def add_two_column_slide(
    prs, title, left_title, left_bullets, right_title, right_bullets
):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Left column title
    left_title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(6), Inches(0.5)
    )
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.1), Inches(6), Inches(0.5)
    )
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    return slide


# ============ SLIDE 1: Title ============
add_title_slide(
    prs,
    "JPX Tokyo Stock Exchange Prediction",
    "Machine Learning for Stock Return Forecasting",
)

# ============ SLIDE 2: Background ============
add_content_slide(
    prs,
    "1. Background",
    [
        "JPX (Japan Exchange Group) is one of the world's largest stock exchanges",
        "Goal: Predict stock returns and build investment portfolios",
        "Traditional analysis cannot meet the needs of complex financial markets",
        "Machine learning provides new approaches for stock prediction",
    ],
)

# ============ SLIDE 3: Data Overview ============
add_content_slide(
    prs,
    "2. Data Overview",
    [
        "stock_prices.csv: ~2.33M records (2017-2021), daily OHLCV data",
        "stock_list.csv: 4,417 stocks with market cap and sector info",
        "financials.csv: ~100K records, quarterly financial statements",
        "trades.csv: ~5K records, weekly investor trading data",
        "options.csv: Options market data including implied volatility",
        "Prediction target: ~2,000 stocks (Universe0)",
    ],
)

# ============ SLIDE 4: Task Objective ============
add_content_slide(
    prs,
    "3. Task Objective",
    [
        "Original (Kaggle): Predict next-day returns, rank stocks",
        "Our Adaptation: Predict 30-day forward returns (reduce noise)",
        "Select Top 200 (long) + Bottom 200 (short) stocks",
        "Evaluate using Sharpe Ratio of daily spread returns",
    ],
)

# ============ SLIDE 5: Why Top/Bottom 200? ============
add_content_slide(
    prs,
    "Why Top 200 & Bottom 200?",
    [
        "Alpha comes from relative value differences, not absolute returns",
        "200 provides sufficient diversification with manageable trading costs",
        "Long-Short strategy hedges market systematic risk (Beta)",
        "Predicting 'who ranks higher' is easier than 'who will rise'",
        "Fixed number ensures stable strategy capacity for comparison",
    ],
)

# ============ SLIDE 6: Feature Engineering ============
add_two_column_slide(
    prs,
    "4. Feature Engineering",
    "Single Data Source",
    [
        "Price returns (1/5/10/20 days)",
        "Volatility (rolling std)",
        "Moving average deviation",
        "Volume change",
        "OHLCV features",
    ],
    "Multi-Source Fusion (Recommended)",
    [
        "Technical indicators",
        "Market cap & sector factors",
        "Financial factors (ROE, P/E)",
        "Investor sentiment (foreign/individual)",
        "Options implied volatility",
    ],
)

# ============ SLIDE 7: Model Selection ============
add_two_column_slide(
    prs,
    "5. Model Selection",
    "Deep Learning (LSTM/Transformer)",
    [
        "Capture long-term dependencies",
        "High training cost",
        "Prone to overfitting",
        "Poor handling of non-stationary data",
    ],
    "Traditional ML (Recommended)",
    [
        "LightGBM: Fast, interpretable, robust",
        "XGBoost: Strong generalization",
        "Ridge: Simple and stable",
        "Our choice: LightGBM + Classification hybrid",
    ],
)

# ============ SLIDE 8: Ensemble Learning ============
add_content_slide(
    prs,
    "6. Ensemble Learning",
    [
        "Random Forest: Same model, multiple instances (bagging)",
        "Multi-Model Ensemble: Different algorithms/data sources combined",
        "Our approach: Weighted ensemble by validation performance",
        "Train separate models for different data sources",
        "Weight = Spearman correlation on validation set",
    ],
)

# ============ SLIDE 9: Data Split & Hyperparameters ============
add_two_column_slide(
    prs,
    "7. Validation & Optimization",
    "Data Split Strategy",
    [
        "Expanding Window: All historical data for training",
        "Rolling Window: Recent 2 years only",
        "Time-series split (no future data leakage)",
        "Validation pairs: (2017→2018), (2018→2019), (2019→2020)",
    ],
    "Hyperparameter Tuning",
    [
        "Extended Grid Search: 50+ parameter combinations",
        "Parameters: n_estimators, learning_rate, max_depth, etc.",
        "Multi-validation pairs for robust selection",
        "Avoid regime change bias",
    ],
)

# ============ SLIDE 10: Evaluation Metrics ============
add_content_slide(
    prs,
    "8. Evaluation Metrics",
    [
        "Spearman Correlation: Measures ranking ability (core metric)",
        "Sharpe Ratio: Risk-adjusted return (portfolio-level)",
        "RMSE: Prediction accuracy",
        "Hit Ratio: Direction accuracy",
        "Priority: Sharpe > Spearman > Hit > RMSE",
    ],
)

# ============ SLIDE 11: Challenges ============
add_content_slide(
    prs,
    "9. Challenges & Difficulties",
    [
        "Historical window size: More data ≠ better prediction",
        "Market regime changes: 2017-2021 had multiple major events",
        "Overfitting risk: Models may memorize past patterns",
        "Non-stationarity: Statistical properties change over time",
    ],
)

# ============ SLIDE 12: Summary ============
add_content_slide(
    prs,
    "10. Summary & Expected Outcomes",
    [
        "Multi-source data integration (prices, financials, trades, options)",
        "LightGBM hybrid model with weighted ensemble",
        "Extended Grid Search for hyperparameter optimization",
        "Expected: Complete prediction pipeline + interpretable strategy",
        "Academic value: ML application in financial prediction",
    ],
)

# ============ SLIDE 13: Thank You ============
add_title_slide(prs, "Thank You!", "Questions & Discussion")

# Save
output_path = "JPX_Stock_Prediction_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
