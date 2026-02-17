"""
将 presentation.html 转换为 PowerPoint 格式
需要安装: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 读取 HTML 内容
with open('presentation.html', 'r', encoding='utf-8') as f:
    html_content = f.read()

# 提取幻灯片内容 (简化解析)
slides_content = []

# 解析每个幻灯片的标题和内容
import re

# 提取所有 slide div
slide_pattern = r'<div class="slide"[^>]*id="slide(\d+)"[^>]*>(.*?)</div>'
slides = re.findall(slide_pattern, html_content, re.DOTALL)

for slide_id, content in slides:
    # 提取标题
    title_match = re.search(r'<h2>(.*?)</h2>', content, re.DOTALL)
    title = title_match.group(1) if title_match else f"Slide {slide_id}"

    # 提取主要文本内容
    text_content = re.sub(r'<[^>]+>', ' ', content)  # 移除 HTML 标签
    text_content = ' '.join(text_content.split())  # 规范化空白

    slides_content.append({
        'id': slide_id,
        'title': title,
        'content': text_content[:500] if len(text_content) > 500 else text_content  # 限制长度
    })

# 创建 PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 设置背景色
background = prs.slide_master.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(26, 26, 46)  # 深蓝色背景

# 标题页样式
title_slide_layout = prs.slide_layouts[6]  # 空白布局
slide = prs.slides.add_slide(title_slide_layout)

# 添加背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(26, 26, 46)
shape.line.fill.background()

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "JPX东京股票交易所预测"
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(0, 212, 255)
p.alignment = PP_ALIGN.CENTER

# 副标题
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "基于机器学习的股票收益率预测与排名"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(170, 170, 170)
p.alignment = PP_ALIGN.CENTER

# 内容页样式
content_layout = prs.slide_layouts[6]

for i, slide_data in enumerate(slides_content[1:], 1):  # 跳过标题页
    slide = prs.slides.add_slide(content_layout)

    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(26, 26, 46)
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 212, 255)
    p.font.bold = True

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # 分段添加内容
    content = slide_data['content']
    paragraphs = content.split('  ')  # 按双空格分段

    for j, para in enumerate(paragraphs[:10]):  # 最多10段
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = para.strip()
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(238, 238, 238)
        p.space_after = Pt(8)

# 保存
output_path = 'presentation.pptx'
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"总页数: {len(prs.slides)}")
